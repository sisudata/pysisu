from pysisu import PySisu
from pysisu.formats import LatestAnalysisResultsFormats
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
import os

# Sisu variables
API_KEY = os.environ.get("SISU_API_KEY")
ANALYSIS_ID = 165576
sisu = PySisu(API_KEY)

# PPTX variables
LAYOUT_TITLE = 0
LAYOUT_TITLE_AND_CONTENT = 1
LAYOUT_TWO_CONTENT = 3
LAYOUT_COMPARISON = 4
LAYOUT_CONTENT_AND_CAPTION = 7
FILENAME = "Sisu Facts.pptx"

# Get facts from Sisu
sisu_table = sisu.get_results(ANALYSIS_ID, confidence_gte="HIGH")
sisu_summary = sisu.get_results(
    ANALYSIS_ID,
    confidence_gte="HIGH",
    format=LatestAnalysisResultsFormats.PROTO,
)
sisu_analyses = sisu.analyses()
sisu_metrics = sisu.metrics()
print("Facts loaded")

# Get general information from analysis
ANALYSIS_NAME = ""
METRIC_NAME = ""
METRIC_ID = ""

for analysis in sisu_analyses.analyses:
    if analysis.id == ANALYSIS_ID:
        ANALYSIS_NAME = analysis.name
        METRIC_ID = analysis.metric_id

# Get the metric name for this analysis
for metric in sisu_metrics.metrics:
    if metric.id == METRIC_ID:
        METRIC_NAME = metric.name

# Create the presentation
p = Presentation()

# Add the overall analysis information to the slide title
sl = p.slide_layouts[LAYOUT_TITLE]
s = p.slides.add_slide(sl)
s.shapes.title.text = ANALYSIS_NAME

sl = p.slide_layouts[LAYOUT_COMPARISON]
s = p.slides.add_slide(sl)
s.shapes.title.text = ANALYSIS_NAME

# Add group names for comparison
s.shapes[1].text = (
    "Group A"
    + "\n"
    + sisu_summary.analysis_result.key_driver_analysis_result.group_comparison.group_a.name
)
s.shapes[3].text = (
    "Group B"
    + "\n"
    + sisu_summary.analysis_result.key_driver_analysis_result.group_comparison.group_b.name
)

# Add metric value, min, max, median, average, sum, rows summary information
summary_card = (
    "Metric Value: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_a_card.summary_value
    )
    + "\n"
)
summary_card += (
    "Min: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_a_card.min
    )
    + "\n"
)
summary_card += (
    "Max: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_a_card.max
    )
    + "\n"
)
summary_card += (
    "Median: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_a_card.median
    )
    + "\n"
)
summary_card += (
    "Average: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_a_card.average
    )
    + "\n"
)
summary_card += (
    "Sum: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_a_card.sum
    )
    + "\n"
)
summary_card += (
    "Rows: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_a_card.total_size
    )
    + "\n"
)
s.shapes[2].text = summary_card

summary_card = (
    "Metric Value: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_b_card.summary_value
    )
    + "\n"
)
summary_card += (
    "Min: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_b_card.min
    )
    + "\n"
)
summary_card += (
    "Max: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_b_card.max
    )
    + "\n"
)
summary_card += (
    "Median: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_b_card.median
    )
    + "\n"
)
summary_card += (
    "Average: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_b_card.average
    )
    + "\n"
)
summary_card += (
    "Sum: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_b_card.sum
    )
    + "\n"
)
summary_card += (
    "Rows: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.group_comparison_card.group_b_card.total_size
    )
    + "\n"
)
s.shapes[4].text = summary_card

# Print facts to the terminal and insert into the presentation deck
print(", ".join([x.column_name for x in sisu_table.header]))

for fact_row in sisu_table.rows:
    pg_row = (
        fact_row.subgroup_id,
        fact_row.confidence,
        fact_row.factor_0_dimension,
        fact_row.factor_0_value,
        fact_row.factor_1_dimension,
        fact_row.factor_1_value,
        fact_row.factor_2_dimension,
        fact_row.factor_2_value,
        fact_row.impact,
        fact_row.group_a_size,
        fact_row.group_b_size,
        fact_row.group_a_value,
        fact_row.group_b_value,
        fact_row.group_a_name,
        fact_row.group_b_name,
    )
    print(pg_row)

    sl = p.slide_layouts[LAYOUT_TWO_CONTENT]
    s = p.slides.add_slide(sl)
    s.shapes.title.text = fact_row.factor_0_dimension
    s.shapes[1].text = (
        "Where "
        + fact_row.factor_0_dimension
        + " is '"
        + str(fact_row.factor_0_value)
        + "'"
    )

    if fact_row.factor_1_dimension:
        s.shapes[1].text = (
            s.shapes[1].text
            + " and "
            + fact_row.factor_1_dimension
            + " is '"
            + str(fact_row.factor_1_value)
            + "'"
        )

    if fact_row.factor_2_dimension:
        s.shapes[1].text = (
            s.shapes[1].text
            + " and "
            + fact_row.factor_2_dimension
            + " is '"
            + str(fact_row.factor_2_value)
            + "'"
        )

    factor_change_val = round(
        fact_row.group_b_value - fact_row.group_a_value, 1
    )
    factor_change_pct = round(
        (
            (fact_row.group_b_value - fact_row.group_a_value)
            / fact_row.group_b_value
        )
        * 100,
        1,
    )

    if factor_change_pct < 0:
        factor_change_dir = "decreased"
    else:
        factor_change_dir = "increased"

    s.shapes[1].text = (
        s.shapes[1].text
        + ", Sisu found with "
        + str(fact_row.confidence).replace("CONFIDENCE_LEVEL_", "")
        + " confidence that the average "
        + METRIC_NAME
        + " "
        + factor_change_dir
        + " "
        + str(abs(factor_change_pct))
        + "% (from "
        + str(round(fact_row.group_a_value, 1))
        + " to "
        + str(round(fact_row.group_b_value, 1))
        + ".) This "
        + factor_change_dir
        + " the overall average "
        + METRIC_NAME
        + " between groups by "
        + str(round(fact_row.impact, 1))
        + "."
    )

    cd = ChartData()
    cd.categories = [METRIC_NAME]
    cd.add_series(fact_row.group_a_name, [round(fact_row.group_a_value, 1)])
    cd.add_series(fact_row.group_b_name, [round(fact_row.group_b_value, 1)])

    x, y, w, h = Inches(5.5), Inches(1.5), Inches(3.5), Inches(5)
    c = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd
    ).chart
    c.has_legend = True
    c.value_axis.minimum_scale = 0
    c.plots[0].has_data_labels = True
    c.plots[0].data_labels.show_value = True
    c.series[0].format.fill.solid()
    c.series[1].format.fill.solid()
    c.series[0].format.fill.fore_color.rgb = RGBColor(34, 100, 245)
    c.series[1].format.fill.fore_color.rgb = RGBColor(252, 196, 13)

    print("Slide inserted successfully into deck")

# Clean up
p.save(FILENAME)
