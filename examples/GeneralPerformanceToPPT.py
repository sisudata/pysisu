from pysisu import PySisu
from pysisu.formats import LatestAnalysisResultsFormats
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
import os

# Sisu variables
API_KEY = os.environ.get("SISU_API_KEY")
ANALYSIS_ID = 165569
sisu = PySisu(API_KEY)

# PPTX variables
LAYOUT_TITLE = 0
LAYOUT_TITLE_AND_CONTENT = 1
LAYOUT_TWO_CONTENT = 3
LAYOUT_COMPARISON = 4
LAYOUT_CONTENT_AND_CAPTION = 7
FILENAME = "Sisu Facts.pptx"

# Get facts from Sisu
sisu_table = sisu.get_results(ANALYSIS_ID, confidence_gte="HIGH")
sisu_summary = sisu.get_results(
    ANALYSIS_ID,
    confidence_gte="HIGH",
    format=LatestAnalysisResultsFormats.PROTO,
)
sisu_analyses = sisu.analyses()
sisu_metrics = sisu.metrics()
print("Facts loaded")

# Get general information from analysis
ANALYSIS_NAME = ""
METRIC_NAME = ""
METRIC_ID = ""

for analysis in sisu_analyses.analyses:
    if analysis.id == ANALYSIS_ID:
        ANALYSIS_NAME = analysis.name
        METRIC_ID = analysis.metric_id

# Get the metric name for this analysis
for metric in sisu_metrics.metrics:
    if metric.id == METRIC_ID:
        METRIC_NAME = metric.name

# Create the presentation
p = Presentation()

# Add the overall analysis information to the slide title
sl = p.slide_layouts[LAYOUT_TITLE]
s = p.slides.add_slide(sl)
s.shapes.title.text = ANALYSIS_NAME

sl = p.slide_layouts[LAYOUT_COMPARISON]
s = p.slides.add_slide(sl)
s.shapes.title.text = ANALYSIS_NAME

s.shapes[1].text = METRIC_NAME

# Add metric value, min, max, median, average, sum, rows summary information
summary_card = (
    "Metric Value: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.general_performance_card.summary_value
    )
    + "\n"
)
summary_card += (
    "Min: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.general_performance_card.min
    )
    + "\n"
)
summary_card += (
    "Max: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.general_performance_card.max
    )
    + "\n"
)
summary_card += (
    "Median: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.general_performance_card.median
    )
    + "\n"
)
summary_card += (
    "Average: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.general_performance_card.average
    )
    + "\n"
)
summary_card += (
    "Sum: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.general_performance_card.sum
    )
    + "\n"
)
summary_card += (
    "Rows: "
    + str(
        sisu_summary.analysis_result.key_driver_analysis_result.summary_card.general_performance_card.total_size
    )
    + "\n"
)
s.shapes[2].text = summary_card

# Print facts to the terminal and insert into the presentation deck
print(", ".join([x.column_name for x in sisu_table.header]))

for fact_row in sisu_table.rows:
    pg_row = (
        fact_row.subgroup_id,
        fact_row.confidence,
        fact_row.factor_0_dimension,
        fact_row.factor_0_value,
        fact_row.factor_1_dimension,
        fact_row.factor_1_value,
        fact_row.factor_2_dimension,
        fact_row.factor_2_value,
        fact_row.impact,
        fact_row.size,
        fact_row.value,
    )
    print(pg_row)

    sl = p.slide_layouts[LAYOUT_TWO_CONTENT]
    s = p.slides.add_slide(sl)
    s.shapes.title.text = fact_row.factor_0_dimension
    s.shapes[1].text = (
        "Where "
        + fact_row.factor_0_dimension
        + " is '"
        + str(fact_row.factor_0_value)
        + "'"
    )

    if fact_row.factor_1_dimension:
        s.shapes[1].text = (
            s.shapes[1].text
            + " and "
            + fact_row.factor_1_dimension
            + " is '"
            + str(fact_row.factor_1_value)
            + "'"
        )

    if fact_row.factor_2_dimension:
        s.shapes[1].text = (
            s.shapes[1].text
            + " and "
            + fact_row.factor_2_dimension
            + " is '"
            + str(fact_row.factor_2_value)
            + "'"
        )

    s.shapes[1].text = (
        s.shapes[1].text
        + ",  with "
        + str(fact_row.confidence).replace("CONFIDENCE_LEVEL_", "")
        + " confidence, the average "
        + METRIC_NAME
        + " value is "
        + str(round(fact_row.value, 1))
        + " with an impact of "
        + str(round(fact_row.impact, 1))
        + "."
    )

    cd = ChartData()
    cd.categories = [METRIC_NAME]
    cd.add_series("Grp Avg", [round(fact_row.value, 1)])

    x, y, w, h = Inches(5.5), Inches(1.5), Inches(3.5), Inches(5)
    c = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd
    ).chart
    c.has_legend = True
    c.value_axis.minimum_scale = 0
    c.plots[0].has_data_labels = True
    c.plots[0].data_labels.show_value = True
    c.series[0].format.fill.solid()
    c.series[0].format.fill.fore_color.rgb = RGBColor(34, 100, 245)

    print("Slide inserted successfully into deck")

# Clean up
p.save(FILENAME)
