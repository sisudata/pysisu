from pysisu import PySisu
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
import os

# Sisu variables
API_KEY = os.environ.get('SISU_API_KEY')
ANALYSIS_ID = 165193
sisu = PySisu(API_KEY)

# PPTX variables
LAYOUT_TITLE = 0
LAYOUT_TITLE_AND_CONTENT = 1
LAYOUT_TWO_CONTENT = 3
LAYOUT_COMPARISON = 4
LAYOUT_CONTENT_AND_CAPTION = 7
FILENAME = 'Sisu Facts.pptx'

# Get facts from Sisu
sisu_table = sisu.get_results(ANALYSIS_ID, {"top_drivers": "True"})
print("Facts loaded")

# Get general information from analysis
ANALYSIS_NAME = 'My Sisu Analysis'  # TO-DO: Add the analysis name here
METRIC_NAME = 'My Metric'  # TO-DO: Add the metric name here

# Create the presentation
p = Presentation()

# Add the overall analysis information to the slide title
sl = p.slide_layouts[LAYOUT_TITLE]
s = p.slides.add_slide(sl)
s.shapes.title.text = ANALYSIS_NAME

sl = p.slide_layouts[LAYOUT_COMPARISON]
s = p.slides.add_slide(sl)
s.shapes.title.text = ANALYSIS_NAME

# TO-DO: add date ranges for period
s.shapes[1].text = 'Previous Period' + '\n' + '<date range>'
s.shapes[3].text = 'Recent Period' + '\n' + '<date range>'

# TO-DO: add metric value, min, max, median, average, sum, rows
s.shapes[2].text = '<previous period summary data>'
s.shapes[4].text = '<recent period summary data>'

# Print facts to the terminal and insert into the presentation deck
print(', '.join([x.column_name for x in sisu_table.header]))

for fact_row in sisu_table.rows:
    pg_row = (fact_row.subgroup_id, fact_row.is_top_driver, fact_row.factor_0_dimension, fact_row.factor_0_value, fact_row.factor_1_dimension, fact_row.factor_1_value, fact_row.factor_2_dimension, fact_row.factor_2_value, fact_row.previous_period_size, fact_row.recent_period_size, fact_row.previous_period_value, fact_row.recent_period_value, fact_row.previous_period_start, fact_row.previous_period_end, fact_row.recent_period_start, fact_row.recent_period_end)
    print(pg_row)

    sl = p.slide_layouts[LAYOUT_TWO_CONTENT]
    s = p.slides.add_slide(sl)
    s.shapes.title.text = fact_row.factor_0_dimension
    s.shapes[1].text = 'Where ' + fact_row.factor_0_dimension + ' is \'' + fact_row.factor_0_value + '\''

    if fact_row.factor_1_dimension:
        s.shapes[1].text = s.shapes[1].text + ' and ' + fact_row.factor_1_dimension + ' is \'' + fact_row.factor_1_value + '\''

    if fact_row.factor_2_dimension:
        s.shapes[1].text = s.shapes[1].text + ' and ' + fact_row.factor_2_dimension + ' is \'' + fact_row.factor_2_value + '\''

    factor_change_val = round(fact_row.recent_period_value - fact_row.previous_period_value, 1)
    factor_change_pct = round(((fact_row.recent_period_value - fact_row.previous_period_value) / fact_row.recent_period_value) * 100, 1)

    if factor_change_pct < 0:
        factor_change_dir = 'decreased'
    else:
        factor_change_dir = 'increased'

    # TO-DO: Add impact score change when supported by the API
    s.shapes[1].text = s.shapes[1].text + ', Sisu found that the average ' + METRIC_NAME + ' ' + factor_change_dir + ' ' + str(abs(factor_change_pct)) + '% (from ' + str(round(fact_row.previous_period_value, 1)) + ' to ' + str(round(fact_row.recent_period_value, 1)) + '.) This ' + factor_change_dir + ' the overall average ' + METRIC_NAME + ' over the same period by ' + '<impact score change>.'

    cd = ChartData()
    cd.categories = [METRIC_NAME]
    cd.add_series('Previous', [round(fact_row.previous_period_value, 1)])
    cd.add_series('Recent', [round(fact_row.recent_period_value, 1)])

    x, y, w, h = Inches(5.5), Inches(1.5), Inches(3.5), Inches(5)
    c = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd).chart
    c.has_legend = True
    c.value_axis.minimum_scale = 0
    c.plots[0].has_data_labels = True
    c.plots[0].data_labels.show_value = True
    c.series[0].format.fill.solid()
    c.series[1].format.fill.solid()
    c.series[0].format.fill.fore_color.rgb = RGBColor(34, 100, 245)
    c.series[1].format.fill.fore_color.rgb = RGBColor(252, 196, 13)

    print("Slide inserted successfully into deck")

# Clean up
p.save(FILENAME)
